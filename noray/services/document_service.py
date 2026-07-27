import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from noray.rag.chunker import select_and_chunk
from noray.rag.embeddings import EmbeddingsManager
from noray.rag.sparse_index import SparseBM25Index
from noray.rag.vector_store import VectorStoreFactory


class DocumentService:
    """Ingestion pipeline that extracts text from documents, chunks them, embeds them, and updates indexes."""
    def __init__(self, vector_store=None, embedder=None, sparse_index=None):
        self.vector_store = vector_store or VectorStoreFactory.get_vector_store()
        self.embedder = embedder or EmbeddingsManager.get_embedder()
        self.sparse_index = sparse_index or SparseBM25Index()

    def parse_file(self, file_path: Path) -> str:
        """Parse documents based on extension, with OCR fallback for images."""
        ext = file_path.suffix.lower()
        if ext == ".pdf":
            return self._parse_pdf(file_path)
        elif ext == ".docx":
            return self._parse_docx(file_path)
        elif ext in [".pptx"]:
            return self._parse_pptx(file_path)
        elif ext in [".xls", ".xlsx"]:
            return self._parse_excel(file_path)
        elif ext in [".md", ".markdown", ".txt", ".csv"]:
            return file_path.read_text(encoding="utf-8", errors="ignore")
        elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
            return self._parse_ocr(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def _parse_pdf(self, file_path: Path) -> str:
        text = ""
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except ImportError:
            # Fallback to PyMuPDF if pdfplumber is not available
            import fitz
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text() + "\n"
        return text

    def _parse_docx(self, file_path: Path) -> str:
        import docx
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    def _parse_pptx(self, file_path: Path) -> str:
        text = ""
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        except Exception as e:
            text = f"[PPTX parsing fallback] Error: {e}\n"
        return text

    def _parse_excel(self, file_path: Path) -> str:
        text = ""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text += f"Sheet: {sheet}\n"
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join([str(val) for val in row if val is not None])
                    if row_text:
                        text += row_text + "\n"
        except Exception as e:
            text = f"[Excel parsing fallback] Error: {e}\n"
        return text

    def _parse_ocr(self, file_path: Path) -> str:
        try:
            import pytesseract
            from PIL import Image
            return pytesseract.image_to_string(Image.open(file_path))
        except Exception as e:
            raise RuntimeError(f"OCR parsing failed: {e}. Check if Tesseract OCR is installed.") from e

    def auto_classify(self, filename: str, text: str) -> dict[str, Any]:
        """Auto-classify document type, language, reading time, keywords, and AI summary."""
        ext = Path(filename).suffix.lower()
        text_lower = text.lower()

        # Classify Document Type
        if "resume" in filename.lower() or "cv" in filename.lower() or "education" in text_lower and "experience" in text_lower:
            doc_type = "Resume"
        elif "scholarship" in text_lower or "grant" in text_lower or "fellowship" in text_lower:
            doc_type = "Scholarship"
        elif "abstract" in text_lower and ("introduction" in text_lower or "references" in text_lower):
            doc_type = "Research Paper"
        elif ext in [".py", ".ts", ".js", ".java", ".cpp", ".c", ".h", ".rs", ".go"]:
            doc_type = "Source Code"
        elif ext in [".xlsx", ".xls", ".csv"]:
            doc_type = "Spreadsheet"
        elif ext in [".pptx"]:
            doc_type = "Presentation"
        else:
            doc_type = "Document"

        words = text.split()
        word_count = len(words)
        reading_time = max(1, round(word_count / 200))

        # Generate summary stub
        summary_text = text[:280].strip().replace("\n", " ") + ("..." if len(text) > 280 else "")

        # Extract top keywords
        common = {"the", "and", "a", "of", "to", "in", "is", "for", "that", "with", "this", "on", "are", "as"}
        clean_words = [w.strip(".,;:()[]\"'").lower() for w in words if len(w) > 3]
        freq = {}
        for w in clean_words:
            if w not in common:
                freq[w] = freq.get(w, 0) + 1
        keywords = sorted(freq.keys(), key=lambda k: freq[k], reverse=True)[:5]

        return {
            "doc_type": doc_type,
            "summary": summary_text,
            "keywords": keywords,
            "language": "en",
            "reading_time_min": reading_time,
            "word_count": word_count,
        }

    def ingest_document(self, file_path: Path, category: str = "general", metadata_overrides: dict[str, Any] | None = None) -> dict[str, Any]:
        """Runs the complete end-to-end extraction, chunking, embedding, and dual-indexing pipeline."""
        filename = file_path.name
        text = self.parse_file(file_path)

        if not text.strip():
            raise ValueError(f"Extracted document text is empty: {filename}")

        # Choose chunking strategy and chunk text
        chunks_data = select_and_chunk(text, filename, embedding_model=self.embedder)

        # Auto-classify and generate metadata
        auto_meta = self.auto_classify(filename, text)

        # Build document-level metadata
        doc_metadata = {
            "source": filename,
            "category": category,
            "doc_type": auto_meta["doc_type"],
            "summary": auto_meta["summary"],
            "keywords": auto_meta["keywords"],
            "language": auto_meta["language"],
            "reading_time_min": auto_meta["reading_time_min"],
            "word_count": auto_meta["word_count"],
            "chunks_count": len(chunks_data),
            "created_at": datetime.now(timezone.utc).isoformat(),
            "updated_at": datetime.now(timezone.utc).isoformat(),
            "version": "1.0.0"
        }
        if metadata_overrides:
            doc_metadata.update(metadata_overrides)

        # 1. Store in Qdrant vector index
        collection_name = "user_documents"

        # Test vector dimension to initialize collection config
        sample_vector = self.embedder.embed(["Hello"])[0]
        vector_dim = len(sample_vector)
        self.vector_store.create_collection(collection_name, vector_dim)

        # Embed all chunks
        chunk_texts = [c["content"] for c in chunks_data]
        embeddings = self.embedder.embed(chunk_texts)

        upsert_points = []
        sparse_chunks = []

        for idx, chunk in enumerate(chunks_data):
            chunk_id = str(uuid.uuid4())
            vector = embeddings[idx]

            # Combine chunk-level details with doc-level metadata
            payload = doc_metadata.copy()
            payload["content"] = chunk["content"]
            payload["chunk_index"] = chunk["chunk_index"]
            payload["strategy"] = chunk["strategy"]

            upsert_points.append({
                "id": chunk_id,
                "vector": vector,
                "payload": payload
            })

            sparse_chunks.append({
                "id": chunk_id,
                "content": chunk["content"],
                "payload": payload
            })

        # Vector upload
        self.vector_store.upsert(collection_name, upsert_points)

        # 2. Store in Sparse index
        # Load existing index, append new chunks, and refit
        existing_chunks = []
        if self.sparse_index.load():
            existing_chunks = self.sparse_index.chunks

        all_chunks = existing_chunks + sparse_chunks
        self.sparse_index.fit_and_save(all_chunks)

        return {
            "filename": filename,
            "chunks_count": len(chunks_data),
            "strategy": chunks_data[0]["strategy"] if chunks_data else "none",
            "category": category
        }
